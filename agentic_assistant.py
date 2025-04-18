#!/usr/bin/env python
# coding: utf-8

import os
from typing import List, Dict, Any, Optional
from langchain.chat_models import ChatOpenAI
from langchain.schema import Document
from langchain.agents import initialize_agent, Tool
from langchain.agents import AgentType
from langchain.chains import RetrievalQA
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
import matplotlib.pyplot as plt
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import requests
from io import BytesIO
from image_creator import DalleImageGenerator
import traceback
import re
from pptx.enum.text import MSO_AUTO_SIZE
import time

class AgenticAssistant:
    def __init__(self, vector_db, model_name="chatgpt-4o-latest", temperature=0.7, api_key=os.getenv("OPENAI_API_KEY")):
        """Initialize the AI-powered assistant with enhanced capabilities."""
        self.vector_db = vector_db
        self.model_name = model_name
        self.temperature = temperature
        self.api_key = api_key
        self.image_generator = DalleImageGenerator()

        # Initialize the language model with enhanced configuration
        self.llm = ChatOpenAI(
            model_name=model_name,
            temperature=temperature,
            max_tokens=4000
        )

        # Initialize conversation memory with larger buffer
        self.memory = ConversationBufferMemory(
            memory_key="chat_history",
            return_messages=True,
            max_len=20
        )

        # Initialize tools with improved descriptions
        self.tools = self._initialize_tools()

        # Initialize the agent with better configuration
        self.agent = initialize_agent(
            self.tools,
            self.llm,
            agent=AgentType.CHAT_CONVERSATIONAL_REACT_DESCRIPTION,
            verbose=True,
            memory=self.memory,
            handle_parsing_errors=True
        )

    def _initialize_tools(self) -> List[Tool]:
        """Initialize and return the tools available to the agent with enhanced descriptions."""
        tools = [
            Tool(
                name="Document Summarization",
                func=self.summarize_document,
                description="Summarizes key points from documents. Input should be a clear query about document content."
            ),
            Tool(
                name="Career Coach",
                func=self.career_coach,
                description="Analyzes job descriptions and resumes. Identifies skill gaps and provides interview prep."
            ),
            Tool(
                name="Chart Analysis",
                func=self.analyze_charts,
                description="Analyzes charts/graphs. Input should be a specific query about visual data."
            ),
            Tool(
                name="Financial Data Insights",
                func=self.analyze_financial_data,
                description="Analyzes financial reports and market trends. Provides investment insights."
            ),
            Tool(
                name="General Question Answering",
                func=self.answer_general_query,
                description="Answers general knowledge questions using available documents."
            ),
            Tool(
                name="STEM Professor - Practice Questions",
                func=self.generate_practice_questions,
                description="Generates 30-50 practice questions from course material."
            ),        
            Tool(
                name="Image Creator",
                func=self.create_single_image,
                description="ONLY for explicit requests to create images. Input MUST be a detailed visual description."
            ),
            Tool(
                name="PowerPoint Generator",
                func=self.generate_presentation,
                description="ONLY for creating PowerPoints. Input should be the presentation topic/subject."
            )
        ]
        return tools

    def run(self, query: str) -> str:
        """
        Enhanced query processing with better error handling and logging.
        """
        print(f"\n--- Processing query: {query[:100]}... ---")
        try:
            response = self.agent.run(query)
            print(f"--- Agent response: {response[:100]}... ---")
            return response
        except Exception as e:
            error_msg = f"ERROR::Agent processing failed: {str(e)}"
            print(f"--- ERROR: {error_msg}\n{traceback.format_exc()} ---")
            return error_msg

    # --- Core Tool Methods ---
    
    def generate_practice_questions(self, query: str) -> str:
        """Enhanced practice question generation with difficulty levels."""
        results = self.vector_db.search(query, k=10)
        if not results:
            return "No relevant course materials found."

        context = "\n\n".join([doc.page_content for doc in results])
        prompt = """
        As a STEM professor, generate 30-50 practice questions across three difficulty levels:
        - Easy (20%)
        - Medium (50%) 
        - Hard (30%)
        
        Material:
        {context}
        
        Format:
        [Difficulty] Question: [Question]
        Answer: [Answer]
        Explanation: [Brief explanation]
        """
        return self._llm_predict(prompt, context=context)

    def analyze_charts(self, query: str) -> str:
        """Enhanced chart analysis with visualization suggestions."""
        results = self.vector_db.search(query, k=5)
        if not results:
            return "No relevant charts found."

        context = "\n\n".join([doc.page_content for doc in results])
        prompt = """
        Analyze these charts and:
        1. Identify key trends
        2. Note anomalies
        3. Suggest alternative visualizations
        4. Provide actionable insights
        
        Charts:
        {context}
        """
        return self._llm_predict(prompt, context=context)

    def analyze_financial_data(self, query: str) -> str:
        """Comprehensive financial analysis with risk assessment."""
        results = self.vector_db.search(query, k=5)
        if not results:
            return "No relevant financial data found."

        context = "\n\n".join([doc.page_content for doc in results])
        prompt = """
        As a financial analyst, provide:
        1. Key performance indicators
        2. Risk assessment
        3. Growth projections
        4. Investment recommendations
        
        Data:
        {context}
        """
        return self._llm_predict(prompt, context=context)

    def answer_general_query(self, query: str) -> str:
        """Enhanced QA with confidence scoring."""
        results = self.vector_db.search(query, k=5)
        if not results:
            return "No relevant information found."

        context = "\n\n".join([doc.page_content for doc in results])
        prompt = """
        Answer the question and rate your confidence (1-5):
        Question: {query}
        Context: {context}
        
        Answer:
        Confidence: [X/5]
        """
        return self._llm_predict(prompt, context=context, query=query)

    def summarize_document(self, query: str) -> str:
        """Enhanced summarization with key takeaways."""
        results = self.vector_db.search(query, k=5)
        if not results:
            return "No relevant documents found."

        context = "\n\n".join([doc.page_content for doc in results])
        prompt = """
        Create a comprehensive summary with:
        1. Key points
        2. Important quotes
        3. Action items
        4. Related concepts
        
        Document:
        {context}
        """
        return self._llm_predict(prompt, context=context)

    def career_coach(self, query: str) -> str:
        """Enhanced career coaching with skill mapping."""
        results = self.vector_db.search(query, k=5)
        if not results:
            return "No relevant career documents found."

        context = "\n\n".join([doc.page_content for doc in results])
        prompt = """
        As a career coach, provide:
        1. Skill gap analysis
        2. Interview preparation plan
        3. Resume improvement suggestions
        4. Career path recommendations
        
        Materials:
        {context}
        """
        return self._llm_predict(prompt, context=context)

    # --- Image Generation ---
    def create_single_image(self, prompt: str) -> str:
        """Robust single image generation with enhanced error handling."""
        print(f"\n--- Image generation requested: {prompt[:100]}... ---")
        if not self.image_generator:
            return "ERROR::Image generation service unavailable."

        try:
            # Clean and validate prompt
            cleaned_prompt = re.sub(
                r"^(create|generate|make|draw)\s+(an?|the)\s+(image|picture|drawing)\s+(of|about)\s+", 
                "", 
                prompt, 
                flags=re.IGNORECASE
            ).strip() or prompt
            
            # Generate and save image
            image_path = self.image_generator.generate_and_save_image(cleaned_prompt)
            
            if os.path.exists(image_path):
                print(f"--- Image successfully generated: {image_path} ---")
                return f"IMAGE_PATH::{image_path}"
            raise ValueError("Generated image file not found")
            
        except Exception as e:
            error_msg = f"Failed to generate image: {str(e)}"
            print(f"--- ERROR: {error_msg}\n{traceback.format_exc()} ---")
            return f"ERROR::{error_msg}"

    # --- PowerPoint Generation ---
    def generate_presentation(self, topic: str, num_slides: int = 5) -> str:
        """Enhanced PowerPoint generation with automatic image inclusion."""
        print(f"\n--- PowerPoint generation requested: {topic} ({num_slides} slides) ---")
        
        try:
            # Get context from vector DB
            context = self._get_presentation_context(topic)
            
            # Generate slide content
            slide_content = self._generate_slide_content(topic, num_slides, context)
            parsed_slides = self._parse_slides(slide_content)[:num_slides]
            
            # Generate images for slides
            self._generate_slide_images(parsed_slides)
            
            # Create PowerPoint file
            ppt_path = self._create_pptx(topic, parsed_slides)
            return f"PPT_PATH::{ppt_path}"
            
        except Exception as e:
            error_msg = f"Presentation generation failed: {str(e)}"
            print(f"--- ERROR: {error_msg}\n{traceback.format_exc()} ---")
            return f"ERROR::{error_msg}"

    # --- Helper Methods ---
    def _llm_predict(self, template: str, **kwargs) -> str:
        """Standardized LLM prediction with error handling."""
        try:
            prompt = PromptTemplate(template=template, input_variables=list(kwargs.keys()))
            return self.llm.predict(prompt.format(**kwargs))
        except Exception as e:
            return f"ERROR::LLM processing failed: {str(e)}"

    def _get_presentation_context(self, topic: str) -> str:
        """Retrieve relevant context for presentation generation."""
        try:
            results = self.vector_db.search(topic, k=3)
            if results:
                return topic + "\n\nRelevant Context:\n" + "\n\n".join(
                    doc.page_content for doc in results if hasattr(doc, 'page_content')
        except Exception as e:
            print(f"--- WARNING: Context retrieval failed: {str(e)} ---")
        return topic

    def _generate_slide_content(self, topic: str, num_slides: int, context: str) -> str:
        """Generate structured slide content using LLM."""
        prompt = f"""
Create content for a {num_slides}-slide presentation about: "{topic}".
Context: {context}

For EACH slide:
1. Title: [Concise, engaging title]
2. Content: [3-5 bullet points or short paragraph]
3. Image Description: [Detailed visual description]

Format response as:
Slide N Title: [Title]
Slide N Content: [Content] 
Slide N Image Description: [Description]
---"""
        return self.llm.predict(prompt)

    def _parse_slides(self, response: str) -> List[Dict]:
        """Parse generated slide content into structured format."""
        slides = []
        pattern = re.compile(
            r"Slide\s*(?P<num>\d+)\s*Title:(?P<title>.*?)\n"
            r"Slide\s*\1\s*Content:(?P<content>.*?)\n"
            r"Slide\s*\1\s*Image Description:(?P<image_prompt>.*?)"
            r"(?=\nSlide\s*\d+\s*Title:|\Z)",
            re.DOTALL | re.IGNORECASE
        )
        
        for match in pattern.finditer(response):
            slide_data = match.groupdict()
            slides.append({
                'title': slide_data['title'].strip(),
                'content': slide_data['content'].strip(),
                'image_prompt': slide_data['image_prompt'].strip()
            })
            
        return slides

    def _generate_slide_images(self, slides: List[Dict]) -> None:
        """Generate images for slides and add paths to slide data."""
        for slide in slides:
            try:
                image_path = self.image_generator.generate_and_save_image(slide['image_prompt'])
                slide['image_path'] = image_path if os.path.exists(image_path) else None
            except Exception as e:
                print(f"--- WARNING: Image generation failed for slide: {str(e)} ---")
                slide['image_path'] = None

    def _create_pptx(self, topic: str, slides: List[Dict]) -> str:
        """Create PowerPoint file from structured slide data."""
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic.title()
        title_slide.placeholders[1].text = "AI-Generated Presentation"
        
        # Content slides
        for i, slide in enumerate(slides, 1):
            try:
                content_slide = prs.slides.add_slide(prs.slide_layouts[5])
                content_slide.shapes.title.text = slide.get('title', f"Slide {i}")
                
                # Add content
                if 'content' in slide:
                    tf = content_slide.placeholders[1].text_frame
                    tf.text = slide['content']
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                
                # Add image if available
                if slide.get('image_path') and os.path.exists(slide['image_path']):
                    content_slide.placeholders[2].insert_picture(slide['image_path'])
                    
            except Exception as e:
                print(f"--- WARNING: Failed to create slide {i}: {str(e)} ---")
        
        # Save presentation
        safe_topic = re.sub(r'[^\w\s-]', '', topic).strip()
        safe_topic = re.sub(r'[-\s]+', '_', safe_topic)
        ppt_path = f"Presentation_{safe_topic[:50]}_{int(time.time())}.pptx"
        prs.save(ppt_path)
        
        return ppt_path
